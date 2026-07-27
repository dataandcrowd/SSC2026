#!/usr/bin/env python3
"""Update SSC2026_presentation.pptx to the calibrated-model results.

Rewrites the numbers on slides 2, 3, 5 and 6 (all of which came from the
pre-calibration model with 500 agents) and appends two results slides built
from the new figures. Run from the `presentation/` directory:

    python update_presentation.py            # writes SSC2026_presentation.pptx
    python update_presentation.py --check    # dry run, prints what would change

Every number here is traceable to paper_update/numbers.md.
"""
import copy, shutil, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "SSC2026_presentation.pptx"
FIGS = os.path.join("..", "output", "figures")
NAVY = RGBColor(0x0D, 0x21, 0x37)
BG = RGBColor(0xF0, 0xF9, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
DRY = "--check" in sys.argv

# slide (1-based) -> shape id -> replacement lines
EDITS = {
    2: {
        18: ["ABM with 2,500 heterogeneous agents (1 agent = 160 vehicles, "
             "calibrated to observed daily counts), Auckland CBD road network. "
             "Three decision rules, same pricing environment, same LoS measurement."],
    },
    3: {
        17: ["Entry probability declines with fee/VoT ratio",
             "No memory, no adaptation",
             "Standard equilibrium assumption",
             "Result: peak V/C 0.122 → 0.094 (−23%); entry 52% → 40%"],
        25: ["Portfolio of predictors based on recent history",
             "Reads yesterday's congestion, not today's fee",
             "Near-universal entry (91%), so little left to deter",
             "Result: peak V/C unchanged under ToU (−0.2%)"],
        33: ["Learns from own departure time experience",
             "Reward penalises fee paid + congestion encountered",
             "Still improving at day 14 — not yet converged",
             "Result: peak V/C 0.122 → 0.093 (−24%); entry 54% → 34%"],
    },
    5: {
        10: ["Fewer drivers enter at the peak: entry falls from 52% to 40% and "
             "peak V/C by about a quarter. The response is smooth and repeats "
             "day after day."],
        14: ["No charge  0.122"],
        16: ["ToU  0.094"],
        18: ["↓  Steady reduction"],
        24: ["Drivers read yesterday's congestion, not today's price. Almost "
             "everyone still enters (91%), so the charge has almost nothing "
             "left to deter and peak congestion is unchanged."],
        28: ["No charge  0.173"],
        30: ["ToU  0.173"],
        32: ["→  No response"],
        38: ["Drivers learn from experience and were still improving on day 14. "
             "Entry falls from 54% to 34%, peak-hour V/C by about half, and "
             "congestion does not shift to nearby roads."],
        42: ["No charge  0.122"],
        44: ["ToU  0.093"],
        46: ["↓  Largest cut in entries"],
    },
    6: {
        30: ["Q-learning gives the most policy-relevant outcome: the cordon "
             "boundary falls furthest (−37%), with no build-up on the periphery"],
        33: ["El Farol shows what happens when agents track each other rather "
             "than the fee — entry stays at 91% and the charge does nothing"],
        49: ["Endogenous Routing"],
        50: ["Agents keep fixed shortest paths, so the only response is "
             "temporal — rerouting would test diversion around the cordon"],
        54: ["Scale beyond 2,500 agents (≈400k vehicles today); dynamic OD "
             "matrix replacing the fixed assumption"],
    },
}

NEW_SLIDES = [
    {
        "title": "Where the Charge Acts Within the Day",
        "image": "sensitivity_hourly_profile.png",
        "caption": "Hour-of-day inner-cordon V/C, days 8–14, no charge (grey) vs ToU "
                   "(blue); shaded bands are the $6 peak and $4 shoulder fee windows. "
                   "Pay −23% at the morning peak, Learn −49%, Oscillate +2%. The "
                   "whole-day mean falls by as much as the peaks do, so the charge "
                   "deters trips rather than spreading them.",
    },
    {
        "title": "No Displacement onto the Ring Roads",
        "image": "map_redistribution.png",
        "caption": "Change in per-link peak V/C under ToU: blue is a reduction, red an "
                   "increase, dashed outline is the cordon. Under Pay and Learn the "
                   "boundary and approach roads fall together with the interior "
                   "(−24% and −37% at the boundary). Under Oscillate the map is "
                   "mottled and the net change is zero.",
    },
]


def set_lines(shape, lines, label):
    tf = shape.text_frame
    paras = tf.paragraphs
    if len(lines) == len(paras):
        for para, line in zip(paras, lines):
            runs = para.runs
            if not runs:
                continue
            runs[0].text = line
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
    elif len(lines) == 1 and paras:
        runs = paras[0].runs
        if runs:
            runs[0].text = lines[0]
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
        for p in list(paras[1:]):
            p._p.getparent().remove(p._p)
    else:
        raise SystemExit(f"{label}: {len(lines)} lines vs {len(paras)} paragraphs")


def add_figure_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    # match the deck: light background, navy header band, white title
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.64))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = spec["title"]
    r.font.name, r.font.size, r.font.bold = "Calibri", Pt(26), True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    path = os.path.join(FIGS, spec["image"])
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    max_w, max_h = Inches(9.2), Inches(3.05)
    scale = min(max_w / w, max_h / h)
    pw, ph_ = int(w * scale), int(h * scale)
    slide.shapes.add_picture(path, int((prs.slide_width - pw) / 2), Inches(1.15),
                             width=pw, height=ph_)

    cap = slide.shapes.add_textbox(Inches(0.4), Inches(4.35), Inches(9.2), Inches(1.0))
    cap.text_frame.word_wrap = True
    cr = cap.text_frame.paragraphs[0].add_run()
    cr.text = spec["caption"]
    cr.font.name, cr.font.size = "Calibri", Pt(11)
    cr.font.color.rgb = MUTED
    cap.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    return slide


def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[from_idx])
    sldIdLst.insert(to_idx, ids[from_idx])


prs = Presentation(SRC)
for slide_no, edits in EDITS.items():
    slide = prs.slides[slide_no - 1]
    by_id = {sh.shape_id: sh for sh in slide.shapes}
    for shape_id, lines in edits.items():
        shape = by_id.get(shape_id)
        if shape is None or not shape.has_text_frame:
            raise SystemExit(f"slide {slide_no}: no text shape with id {shape_id}")
        old = shape.text_frame.text.replace("\n", " | ")[:60]
        print(f"slide {slide_no} id{shape_id}: {old!r} -> {lines[0][:60]!r}")
        if not DRY:
            set_lines(shape, lines, f"slide {slide_no} id{shape_id}")

if not DRY:
    n_before = len(prs.slides._sldIdLst)
    for i, spec in enumerate(NEW_SLIDES):
        add_figure_slide(prs, spec)
        move_slide(prs, n_before + i, 5 + i)   # after slide 5 (the results slide)
        print(f"added slide: {spec['title']}")
    shutil.copy(SRC, SRC + ".bak")
    prs.save(SRC)
    print(f"saved {SRC} (backup: {SRC}.bak)")
else:
    print("(dry run, nothing written)")
